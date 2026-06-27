#!/usr/bin/env python3
"""
HTML to PowerPoint Converter for Visual Workbench
Converts HTML/CSS presentation screens to PPTX format.
"""

import sys
import json
import base64
from pathlib import Path
from io import BytesIO
from typing import List, Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def create_presentation(screens: List[Dict[str, Any]], output_path: str = "presentation.pptx"):
    """Create a PowerPoint presentation from HTML screens."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, screen in enumerate(screens):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.333)
        height = Inches(6.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        title = screen.get('name', f'Slide {i + 1}')
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        content = extract_text_from_html(screen.get('html', ''))
        if content:
            p2 = tf.add_paragraph()
            p2.text = content[:500]
            p2.font.size = Pt(14)
            p2.alignment = PP_ALIGN.LEFT
        
        slide.shapes.add_textbox(
            Inches(0.25),
            Inches(6.75),
            Inches(2),
            Inches(0.5)
        ).text_frame.paragraphs[0].text = f"Slide {i + 1}"
    
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    return output_path


def extract_text_from_html(html: str) -> str:
    """Extract readable text from HTML content."""
    import re
    text = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL)
    text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    return text[:1000]


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python export-ppt.py <screens.json>")
        print("  Reads HTML screens from stdin if no file provided")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "presentation.pptx"
    
    with open(input_file, 'r') as f:
        screens = json.load(f)
    
    create_presentation(screens, output_file)